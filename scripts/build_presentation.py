import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import os

# Paths
download_dir = os.path.expanduser('~/Downloads')
PDF_PATH = None
for fname in os.listdir(download_dir):
    if 'brainwave' in fname.lower() and fname.lower().endswith('.pdf'):
        PDF_PATH = os.path.join(download_dir, fname)
        break
if not PDF_PATH:
    raise FileNotFoundError('Brainwave PDF template not found in ~/Downloads')

OUT_DIR = os.path.abspath('presentation_assets')
PPTX_OUT = os.path.abspath('CivicShield_AI_Brainwave.pptx')
UPLOADS_DIR = os.path.abspath('backend/app/uploads')

os.makedirs(OUT_DIR, exist_ok=True)

# Convert PDF template pages to images
doc = fitz.open(PDF_PATH)
template_images = []
for i, page in enumerate(doc):
    pix = page.get_pixmap(dpi=150)
    img_path = os.path.join(OUT_DIR, f'template_page_{i+1}.png')
    pix.save(img_path)
    template_images.append(img_path)
print(f'Converted {len(template_images)} template pages')

# Collect some project images from uploads (limit 6)
project_images = []
if os.path.isdir(UPLOADS_DIR):
    for name in os.listdir(UPLOADS_DIR):
        if name.lower().endswith(('.png', '.jpg', '.jpeg')):
            project_images.append(os.path.join(UPLOADS_DIR, name))
project_images = sorted(project_images)[:8]
print(f'Found {len(project_images)} project images')

# Load textual data (use snapshot content extracted earlier if available)
# For accuracy, embed values observed from live site snapshot
stats = {
    'total_reports': '83',
    'critical': '26',
    'high': '18',
    'active_issues': '70',
    'hotspots': '10',
    'avg_priority': '65.8',
    'city_health': '39%',
    'sla_compliance': '38.5%'
}

# Create PPTX and use template page images as backgrounds for slides
prs = Presentation()
# set slide size to default (PowerPoint default 10 x 7.5 inches)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper to add background image
from pptx.enum.shapes import MSO_SHAPE

def add_slide_with_bg(bg_path):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    left = top = Inches(0)
    pic = slide.shapes.add_picture(bg_path, left, top, width=prs.slide_width, height=prs.slide_height)
    return slide

# 1. Title slide - use template page 1 as background
if template_images:
    add_slide_with_bg(template_images[0])
    title_slide = prs.slides[-1]
    # Add project title textbox
    tx = title_slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = 'CivicShield AI'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p = tf.add_paragraph()
    p.text = 'Smart city operations • AI-driven incident detection & response'
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(210,210,210)

# 2. Problem slide
if len(template_images) > 1:
    add_slide_with_bg(template_images[1])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(4.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Problem: Urban infrastructure issues go undetected or take too long to resolve'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p = tf.add_paragraph()
    p.text = 'Citizens report incidents sporadically; authorities lack prioritized, verifiable, and routed reports with SLA tracking.'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(230,230,230)

# 3. Solution slide
if len(template_images) > 2:
    add_slide_with_bg(template_images[2])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(5.5), Inches(4.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = 'Solution: CivicShield AI'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p = tf.add_paragraph()
    p.text = 'AI-powered image analysis (YOLO), priority engine, hotspot detection, SLA manager, and authority dashboard for real-time operations.'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(230,230,230)
    # add screenshot if available
    if project_images:
        s.shapes.add_picture(project_images[0], Inches(6.4), Inches(1.0), width=Inches(3.0))

# 4. Architecture slide
if len(template_images) > 3:
    add_slide_with_bg(template_images[3])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(1.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = 'Architecture'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    tx2 = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(8.6), Inches(4.0))
    tf2 = tx2.text_frame
    tf2.text = 'Frontend (Vite + React) – FastAPI backend – AI services (YOLO inference) – PostgreSQL/SQLite – Authority dashboard & Map APIs.'
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.color.rgb = RGBColor(230,230,230)

# 5. AI / YOLO slide
if len(template_images) > 4:
    add_slide_with_bg(template_images[4])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(5.0), Inches(1.0))
    tf = tx.text_frame
    tf.text = 'AI / YOLO Analysis'
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    # Insert sample detection images
    x = 0.6
    for idx, img in enumerate(project_images[:4]):
        s.shapes.add_picture(img, Inches(0.6 + idx*2.4), Inches(2.0), width=Inches(2.2))

# 6. Priority engine
if len(template_images) > 5:
    add_slide_with_bg(template_images[5])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(2.0))
    tf = tx.text_frame
    tf.text = 'Priority Engine'
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tx2 = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(8.6), Inches(3.0))
    tf2 = tx2.text_frame
    tf2.text = f"Priority score, reasons, and recommended department computed per report. Example metrics:\nTotal reports: {stats['total_reports']}  Critical: {stats['critical']}  High: {stats['high']}\nAvg priority: {stats['avg_priority']}"
    tf2.paragraphs[0].font.size = Pt(14)
    tf2.paragraphs[0].font.color.rgb = RGBColor(230,230,230)

# 7. Hotspot detection
if len(template_images) > 6:
    add_slide_with_bg(template_images[6])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(2.0))
    tf = tx.text_frame
    tf.text = 'Hotspot Detection'
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tx2 = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(8.6), Inches(3.0))
    tf2 = tx2.text_frame
    tf2.text = f"Hotspots detected: {stats['hotspots']}\nHotspot radius and related report counts shown on report details page."
    tf2.paragraphs[0].font.size = Pt(14)

# 8. SLA & Dashboard slide
if len(template_images) > 7:
    add_slide_with_bg(template_images[7])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(2.0))
    tf = tx.text_frame
    tf.text = 'SLA Manager & Dashboard'
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tx2 = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.6), Inches(3.0))
    tf2 = tx2.text_frame
    tf2.text = f"City health: {stats['city_health']}  SLA compliance: {stats['sla_compliance']}\nActive issues: {stats['active_issues']}"
    tf2.paragraphs[0].font.size = Pt(14)
    # attach a dashboard screenshot if available
    if len(project_images) > 1:
        s.shapes.add_picture(project_images[1], Inches(6.4), Inches(2.4), width=Inches(3.0))

# 9. Report tracking and resolution verification
if len(template_images) > 8:
    add_slide_with_bg(template_images[8])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(1.6))
    tf = tx.text_frame
    tf.text = 'Report Tracking & Resolution Verification'
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tx2 = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.6), Inches(3.0))
    tf2 = tx2.text_frame
    tf2.text = 'Public tracking (report ID), timeline, status workflow, and AI resolution verification score are integrated into the authority UI.'
    tf2.paragraphs[0].font.size = Pt(14)
    if len(project_images) > 2:
        s.shapes.add_picture(project_images[2], Inches(6.4), Inches(2.4), width=Inches(3.0))

# 10. Impact slide
if len(template_images) > 9:
    add_slide_with_bg(template_images[9])
    s = prs.slides[-1]
    tx = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.6), Inches(3.0))
    tf = tx.text_frame
    tf.text = 'Impact'
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf2 = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.6), Inches(3.0))
    tf2.text_frame.text = 'Faster detection and routing, measurable SLA improvements, visual resolution verification, and prioritized field deployments.'
    tf2.text_frame.paragraphs[0].font.size = Pt(14)

# Save presentation
prs.save(PPTX_OUT)
print('Saved PPTX to', PPTX_OUT)
