import os
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

# Use existing Brainwave PPTX as template
TEMPLATE_PPTX = 'CivicShield_AI_Brainwave.pptx'
OUTPUT_PPTX = 'CivicShield_AI_Brainwave_FINAL.pptx'
OUTPUT_PDF = 'CivicShield_AI_Brainwave_FINAL.pdf'

# Ensure output directories exist
os.makedirs('presentation_assets', exist_ok=True)

# Use repository uploads as image sources
uploads_dir = os.path.abspath('backend/app/uploads')
if os.path.isdir(uploads_dir):
    upload_images = [os.path.join(uploads_dir, f) for f in os.listdir(uploads_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    upload_images = sorted(upload_images)
else:
    upload_images = []

fallback_images = upload_images[:3]

template_images = []
for i in range(1, 5):
    path = os.path.join('presentation_assets', f'template_page_{i}.png')
    if os.path.exists(path):
        template_images.append(path)

# Final slide content
slides = [
    {
        'title': 'CIVICSHIELD AI',
        'subtitle': 'AI-Powered Smart City Civic Infrastructure Command Center',
        'body': 'Detect • Prioritize • Dispatch • Verify',
        'bullets': None,
        'image': fallback_images[0] if fallback_images else None,
        'bg': template_images[0] if len(template_images) > 0 else None,
    },
    {
        'title': 'THE CIVIC INFRASTRUCTURE CHALLENGE',
        'subtitle': None,
        'body': 'Citizen reports are fragmented, manual prioritization delays response, and verification is missing. CivicShield AI addresses potholes, garbage, broken streetlights, road damage, water leaks and infrastructure failures.',
        'bullets': [
            'Delayed response',
            'Manual prioritization',
            'Scattered reports',
            'Repeated duplicate issues',
            'Poor SLA visibility',
            'Weak resolution verification',
        ],
        'image': None,
        'bg': template_images[1] if len(template_images) > 1 else template_images[0] if template_images else None,
    },
    {
        'title': 'CIVICSHIELD AI',
        'subtitle': 'End-to-End Smart City Command Pipeline',
        'body': 'Citizen / Image → AI Vision → Detection → Severity & Safety → Priority → Hotspot → Department → SLA → Verification',
        'bullets': [
            'AI-powered',
            'Real-time',
            'Priority-driven',
            'Location-aware',
            'Resolution-aware',
        ],
        'image': fallback_images[1] if len(fallback_images) > 1 else None,
        'bg': template_images[2] if len(template_images) > 2 else template_images[0] if template_images else None,
    },
    {
        'title': 'AI-POWERED CIVIC INCIDENT INTELLIGENCE',
        'subtitle': None,
        'body': 'Pothole detected with 95% confidence. Severity 8.8/10, Safety Risk HIGH, Priority 92/100. Recommended department: ROAD MAINTENANCE.',
        'bullets': [
            'AI Detection',
            'Priority Engine',
            'Department Recommendation',
        ],
        'image': fallback_images[2] if len(fallback_images) > 2 else None,
        'bg': template_images[3] if len(template_images) > 3 else template_images[0] if template_images else None,
    },
    {
        'title': 'REAL-TIME CIVIC OPERATIONS',
        'subtitle': None,
        'body': 'Live dashboard monitoring with 83 total reports, 26 critical incidents, 18 high-priority issues, hotspot awareness, SLA tracking and map-driven dispatch.',
        'bullets': [
            '83 Total Reports',
            '26 Critical',
            '18 High Priority',
            'Civic Hotspots',
            'SLA Monitoring',
            'Live Map',
        ],
        'image': fallback_images[0] if fallback_images else None,
        'bg': template_images[0] if template_images else None,
    },
    {
        'title': 'FROM DETECTION TO RESOLUTION',
        'subtitle': None,
        'body': 'Hotspots cluster nearby reports for escalation. SLA intelligence tracks targets, deadlines, due soon and breaches. Resolution verification compares before/after evidence with AI validation.',
        'bullets': [
            'Hotspot cluster detection',
            'Targeted SLA tracking',
            'Evidence-based resolution verification',
        ],
        'image': fallback_images[1] if len(fallback_images) > 1 else None,
        'bg': template_images[1] if len(template_images) > 1 else template_images[0] if template_images else None,
    },
]

# Save the final PPTX
ppt = Presentation(TEMPLATE_PPTX)
for _ in range(len(ppt.slides)):
    rId = ppt.slides._sldIdLst[0].rId
    ppt.part.drop_rel(rId)
    del ppt.slides._sldIdLst[0]

for slide_data in slides:
    slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = slide_data['title']
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    if slide_data['subtitle']:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.8))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = slide_data['subtitle']
        subtitle_tf.paragraphs[0].font.size = Pt(20)
        subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(200,200,220)
    if slide_data.get('body'):
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(5.5), Inches(3.0))
        body_tf = body_box.text_frame
        body_tf.text = slide_data['body']
        body_tf.paragraphs[0].font.size = Pt(16)
        body_tf.paragraphs[0].font.color.rgb = RGBColor(220,220,230)
    if slide_data.get('bullets'):
        bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(5.5), Inches(2.5))
        bullets_tf = bullets_box.text_frame
        for i, bullet in enumerate(slide_data['bullets']):
            if i == 0:
                p = bullets_tf.paragraphs[0]
            else:
                p = bullets_tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200,200,220)
            p.font.bold = False
    if slide_data.get('image') and os.path.exists(slide_data['image']):
        slide.shapes.add_picture(slide_data['image'], Inches(6.2), Inches(2.3), width=Inches(3.0))

ppt.save(OUTPUT_PPTX)
print('Saved', OUTPUT_PPTX)

# Now build the PDF directly from template backgrounds and rendered slide graphics

def load_font(size):
    for font_name in ['Arial.ttf', 'DejaVuSans.ttf', 'LiberationSans-Regular.ttf']:
        try:
            return ImageFont.truetype(font_name, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def wrap_text(text, font, max_width, draw):
    lines = []
    for paragraph in text.split('\n'):
        words = paragraph.split(' ')
        line = ''
        for word in words:
            test = f'{line} {word}'.strip()
            if draw.textlength(test, font=font) <= max_width:
                line = test
            else:
                if line:
                    lines.append(line)
                line = word
        if line:
            lines.append(line)
    return lines


slide_image_paths = []
for idx, slide_data in enumerate(slides, start=1):
    if slide_data['bg'] and os.path.exists(slide_data['bg']):
        slide_img = Image.open(slide_data['bg']).convert('RGB')
    else:
        slide_img = Image.new('RGB', (1600, 1200), color=(18, 24, 52))
    draw = ImageDraw.Draw(slide_img)
    title_font = load_font(56)
    subtitle_font = load_font(28)
    body_font = load_font(24)
    bullet_font = load_font(22)
    accent_color = (255, 255, 255)
    muted_color = (210, 210, 225)
    x_margin = 80
    y = 80
    draw.text((x_margin, y), slide_data['title'], font=title_font, fill=accent_color)
    y += 90
    if slide_data['subtitle']:
        draw.text((x_margin, y), slide_data['subtitle'], font=subtitle_font, fill=muted_color)
        y += 60
    if slide_data['body']:
        max_width = slide_img.width - x_margin * 2
        lines = wrap_text(slide_data['body'], body_font, max_width, draw)
        for line in lines:
            draw.text((x_margin, y), line, font=body_font, fill=muted_color)
            y += 34
        y += 10
    if slide_data.get('bullets'):
        for bullet in slide_data['bullets']:
            bullet_lines = wrap_text(f'• {bullet}', bullet_font, slide_img.width - x_margin * 2, draw)
            for line in bullet_lines:
                draw.text((x_margin, y), line, font=bullet_font, fill=muted_color)
                y += 30
        y += 10
    if slide_data.get('image') and os.path.exists(slide_data['image']):
        try:
            extra = Image.open(slide_data['image']).convert('RGB')
            extra.thumbnail((520, 520), Image.LANCZOS)
            slide_img.paste(extra, (slide_img.width - extra.width - x_margin, 280))
        except OSError:
            pass
    output_path = os.path.join('presentation_assets', f'final_slide_{idx}.png')
    slide_img.save(output_path, format='PNG')
    slide_image_paths.append(output_path)

if slide_image_paths:
    images = [Image.open(p).convert('RGB') for p in slide_image_paths]
    images[0].save(OUTPUT_PDF, save_all=True, append_images=images[1:])
    print('Saved', OUTPUT_PDF)
else:
    print('No slide images were generated for PDF conversion.')
